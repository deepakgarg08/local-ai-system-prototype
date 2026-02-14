from pathlib import Path
from pptx import Presentation
from .base_loader import BaseDocumentLoader


class PptxLoader(BaseDocumentLoader):

    def load(self, path: Path) -> dict:
        prs = Presentation(path)
        text_runs = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)

        return {
            "text": "\n".join(text_runs),
            "metadata": {
                "source": str(path),
                "file_type": "pptx",
                "num_slides": len(prs.slides)
            }
        }
